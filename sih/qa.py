"""Geometric QA for the SIH deck.

There is no LibreOffice on this machine, so the slides cannot be rendered and
inspected. This measures instead: every shape against the slide bounds, every
pair of my own shapes for overlap, and every text frame's wrapped height
against the box it has to fit in - using the real Calibri and Cambria metrics
from the Windows font directory rather than a guess at character width.
"""

from __future__ import annotations

import os
import pathlib

from PIL import ImageFont
from pptx import Presentation
from pptx.util import Emu

DECK = pathlib.Path(__file__).parent / "INGRES-SIH2026-Idea.pptx"
FONTS = pathlib.Path(os.environ["WINDIR"]) / "Fonts"
SLIDE_W, SLIDE_H = 13.333, 7.5
FOOTER_Y = 6.95           # the template's footer bar
PX_PER_PT = 4             # render at 4x for sub-point measurement accuracy

_cache: dict[tuple[str, float], ImageFont.FreeTypeFont] = {}


def font(name: str, size_pt: float):
    key = (name, round(size_pt, 1))
    if key not in _cache:
        f = "cambria.ttc" if name.lower().startswith("cambria") else "calibri.ttf"
        _cache[key] = ImageFont.truetype(str(FONTS / f), int(size_pt * PX_PER_PT))
    return _cache[key]


def inches(v) -> float:
    return Emu(v).inches if v is not None else 0.0


def text_width_in(text: str, name: str, size_pt: float) -> float:
    """Rendered width in inches."""
    f = font(name, size_pt)
    return f.getlength(text) / PX_PER_PT / 72.0


def wrapped_height_in(tf, box_w_in: float) -> float:
    """Height the text frame needs once wrapped into box_w_in."""
    inner = box_w_in - inches(tf.margin_left) - inches(tf.margin_right)
    total = inches(tf.margin_top) + inches(tf.margin_bottom)
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            total += 0.08
            continue
        size = max((r.font.size.pt if r.font.size else 12) for r in runs)
        name = next((r.font.name for r in runs if r.font.name), "Calibri")
        text = "".join(r.text for r in runs)

        # greedy wrap on spaces
        lines, cur = 1, ""
        for word in text.split(" "):
            trial = word if not cur else cur + " " + word
            if text_width_in(trial, name, size) <= inner or not cur:
                cur = trial
            else:
                lines += 1
                cur = word
        before = p.space_before.pt if p.space_before else 0
        after = p.space_after.pt if p.space_after else 0
        total += (lines * size * 1.22 + before + after) / 72.0
    return total


def main() -> int:
    prs = Presentation(str(DECK))
    problems: list[str] = []
    warnings: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        boxes = []
        for sh in slide.shapes:
            x, y = inches(sh.left), inches(sh.top)
            w, h = inches(sh.width), inches(sh.height)
            name = (sh.text_frame.text.strip().split("\n")[0][:38]
                    if sh.has_text_frame and sh.text_frame.text.strip()
                    else str(sh.shape_type))
            boxes.append((name, x, y, w, h, sh))

            # --- bounds ---
            if x < -0.01 or y < -0.7 or x + w > SLIDE_W + 0.01 or y + h > SLIDE_H + 0.01:
                # the template itself hangs a couple of decorations off-canvas
                if sh.shape_id not in (36, 37, 8, 2, 5):
                    problems.append(
                        f"S{i}: '{name}' out of bounds "
                        f"({x:.2f},{y:.2f} {w:.2f}x{h:.2f})")

            # --- text overflow ---
            if sh.has_text_frame and sh.text_frame.text.strip() and w > 0.4:
                need = wrapped_height_in(sh.text_frame, w)
                if need > h + 0.02:
                    over = need - h
                    tag = problems if over > 0.12 else warnings
                    tag.append(
                        f"S{i}: '{name}' text needs {need:.2f}\" in {h:.2f}\" "
                        f"(over by {over:.2f}\")")

            # --- footer collision ---
            if y + h > FOOTER_Y + 0.01 and h < 3 and sh.shape_id not in (9, 10, 6, 7):
                warnings.append(f"S{i}: '{name}' reaches the footer bar "
                                f"(bottom {y + h:.2f}\")")

        # --- overlap between content shapes I added (ignore template chrome) ---
        content = [b for b in boxes if b[3] > 0.3 and b[4] > 0.2 and b[5].shape_id > 100]
        for a in range(len(content)):
            for b in range(a + 1, len(content)):
                n1, x1, y1, w1, h1, _ = content[a]
                n2, x2, y2, w2, h2, _ = content[b]
                ox = min(x1 + w1, x2 + w2) - max(x1, x2)
                oy = min(y1 + h1, y2 + h2) - max(y1, y2)
                if ox <= 0.06 or oy <= 0.06:
                    continue
                # Content sitting inside a background panel is layering, not a
                # collision. Only flag partial overlaps.
                inside = (
                    (x1 <= x2 + 0.02 and y1 <= y2 + 0.02
                     and x1 + w1 >= x2 + w2 - 0.02 and y1 + h1 >= y2 + h2 - 0.02)
                    or (x2 <= x1 + 0.02 and y2 <= y1 + 0.02
                        and x2 + w2 >= x1 + w1 - 0.02 and y2 + h2 >= y1 + h1 - 0.02)
                )
                if inside:
                    continue
                warnings.append(
                    f"S{i}: '{n1}' overlaps '{n2}' by {ox:.2f}x{oy:.2f}\"")

    print(f"slides: {len(prs.slides)}")
    print(f"\nPROBLEMS ({len(problems)})")
    for p in problems:
        print("  x " + p)
    print(f"\nwarnings ({len(warnings)})")
    for w in warnings:
        print("  ! " + w)
    return 1 if problems else 0


if __name__ == "__main__":
    raise SystemExit(main())
