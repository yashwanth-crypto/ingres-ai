"""Render the deck to PNGs so it can actually be looked at.

There is no LibreOffice on this machine, so this draws the slides from the
shape geometry directly: fills, outlines, and text wrapped with the real
Windows font metrics. It is an approximation of PowerPoint's layout engine,
not a substitute for it — but it is enough to see overlaps, overflow,
alignment and colour, which is what visual QA is for.
"""

from __future__ import annotations

import os
import pathlib

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

HERE = pathlib.Path(__file__).parent
DECK = HERE / "INGRES-SIH2026-Idea.pptx"
OUTDIR = HERE / "preview"
FONTS = pathlib.Path(os.environ["WINDIR"]) / "Fonts"

DPI = 110
W, H = int(13.333 * DPI), int(7.5 * DPI)

_fc: dict = {}


def fnt(name: str, pt: float, bold: bool):
    key = (name or "Calibri", round(pt, 1), bold)
    if key in _fc:
        return _fc[key]
    fam = (name or "Calibri").lower()
    if fam.startswith("cambria"):
        f = "cambriab.ttf" if bold else "cambria.ttc"
    else:
        f = "calibrib.ttf" if bold else "calibri.ttf"
    path = FONTS / f
    if not path.exists():
        path = FONTS / ("calibrib.ttf" if bold else "calibri.ttf")
    _fc[key] = ImageFont.truetype(str(path), max(int(pt * DPI / 72), 6))
    return _fc[key]


def px(v):
    return Emu(v).inches * DPI if v is not None else 0


def rgb(color, default=None):
    try:
        if color and color.rgb is not None:
            r, g, b = color.rgb[0], color.rgb[1], color.rgb[2]
            return (r, g, b)
    except Exception:
        pass
    return default


def shape_fill(sh):
    try:
        if sh.fill.type is not None and str(sh.fill.type) != "MSO_FILL_TYPE.BACKGROUND (5)":
            return rgb(sh.fill.fore_color, None)
    except Exception:
        pass
    return None


def draw_shape(d, sh, x, y, w, h):
    fill = shape_fill(sh)
    outline = None
    try:
        if sh.line.fill.type is not None and str(sh.line.fill.type) != "MSO_FILL_TYPE.BACKGROUND (5)":
            outline = rgb(sh.line.color, None)
    except Exception:
        pass
    if fill is None and outline is None:
        return

    name = ""
    try:
        name = str(sh.shape_type)
        auto = sh._element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom")
        if auto is not None:
            name = auto.get("prst") or name
    except Exception:
        pass

    box = [x, y, x + w, y + h]
    if "chevron" in name:
        tip = min(w * 0.18, h * 0.5)
        d.polygon([(x, y), (x + w - tip, y), (x + w, y + h / 2), (x + w - tip, y + h),
                   (x, y + h), (x + tip, y + h / 2)], fill=fill, outline=outline)
    elif "ellipse" in name or "oval" in name:
        d.ellipse(box, fill=fill, outline=outline)
    elif "roundRect" in name:
        d.rounded_rectangle(box, radius=min(w, h) * 0.18, fill=fill, outline=outline)
    elif "downArrow" in name:
        sw = w * 0.42
        d.polygon([(x + (w - sw) / 2, y), (x + (w + sw) / 2, y),
                   (x + (w + sw) / 2, y + h * 0.55), (x + w, y + h * 0.55),
                   (x + w / 2, y + h), (x, y + h * 0.55),
                   (x + (w - sw) / 2, y + h * 0.55)], fill=fill, outline=outline)
    else:
        d.rectangle(box, fill=fill, outline=outline)


def wrap(text, font, max_w):
    out, cur = [], ""
    for word in text.split(" "):
        t = word if not cur else cur + " " + word
        if font.getlength(t) <= max_w or not cur:
            cur = t
        else:
            out.append(cur)
            cur = word
    out.append(cur)
    return out


def draw_text(d, sh, x, y, w, h):
    if not sh.has_text_frame:
        return
    tf = sh.text_frame
    if not tf.text.strip():
        return
    ml, mr = px(tf.margin_left), px(tf.margin_right)
    mt = px(tf.margin_top)
    inner = max(w - ml - mr, 10)

    # total height first, for vertical centring
    blocks = []
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            blocks.append((None, 0, 6, None, None, 0, 0))
            continue
        size = max((r.font.size.pt if r.font.size else 12) for r in runs)
        bold = any(r.font.bold for r in runs)
        fname = next((r.font.name for r in runs if r.font.name), "Calibri")
        col = next((rgb(r.font.color, None) for r in runs if rgb(r.font.color, None)), (30, 26, 23))
        f = fnt(fname, size, bold)
        lines = wrap("".join(r.text for r in runs), f, inner)
        sb = p.space_before.pt * DPI / 72 if p.space_before else 0
        sa = p.space_after.pt * DPI / 72 if p.space_after else 0
        lh = size * 1.22 * DPI / 72
        blocks.append((lines, len(lines) * lh, lh, f, col, sb, sa))
        blocks[-1] = (lines, len(lines) * lh, lh, f, col, sb, sa)
        blocks[-1] += (p.alignment,)
    total = sum(b[1] + b[5] + b[6] for b in blocks if b[0] is not None) + \
            sum(b[2] for b in blocks if b[0] is None)

    anchor = str(tf.vertical_anchor or "")
    cy = y + mt
    if "MIDDLE" in anchor:
        cy = y + max((h - total) / 2, 0)

    for b in blocks:
        if b[0] is None:
            cy += b[2]
            continue
        lines, _, lh, f, col, sb, sa, align = b
        cy += sb
        for ln in lines:
            tx = x + ml
            if align == PP_ALIGN.CENTER:
                tx = x + (w - f.getlength(ln)) / 2
            elif align == PP_ALIGN.RIGHT:
                tx = x + w - mr - f.getlength(ln)
            d.text((tx, cy), ln, font=f, fill=col)
            cy += lh
        cy += sa


def main():
    OUTDIR.mkdir(exist_ok=True)
    for old in OUTDIR.glob("slide-*.png"):
        old.unlink()
    prs = Presentation(str(DECK))
    for i, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (W, H), "white")
        d = ImageDraw.Draw(img)
        for sh in slide.shapes:
            x, y = px(sh.left), px(sh.top)
            w, h = px(sh.width), px(sh.height)
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                d.rectangle([x, y, x + w, y + h], fill=(232, 232, 232))
                d.text((x + 6, y + 6), "[logo]", font=fnt("Calibri", 9, False),
                       fill=(140, 140, 140))
                continue
            draw_shape(d, sh, x, y, w, h)
            draw_text(d, sh, x, y, w, h)
        d.rectangle([0, 0, W - 1, H - 1], outline=(190, 190, 190))
        p = OUTDIR / f"slide-{i}.png"
        img.save(p)
        print(p)


if __name__ == "__main__":
    main()
