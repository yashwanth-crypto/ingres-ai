"""Fill the SIH 2026 idea template with INGRES content and native diagrams.

Run from the repo root:  python sih/build.py

Design notes
------------
Diagrams are built from PowerPoint shapes rather than inserted as images, so
they stay vector-sharp at any zoom, survive the PDF export the portal requires,
and can be nudged by hand afterwards.

The template's "idea details pointers" are reused verbatim as sub-headings -
they are read out of the template rather than retyped, so they cannot drift
from what SIH supplied.
"""

from __future__ import annotations

import copy
import json
import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = pathlib.Path(__file__).parent
TEMPLATE = HERE / "template.pptx"
OUT = HERE / "INGRES-SIH2026-Idea.pptx"

# ---- palette: the project's own, from tailwind.config.js and AquiferHero ----
DEPTH_900 = RGBColor(0x0B, 0x2C, 0x34)
DEPTH_700 = RGBColor(0x16, 0x4E, 0x5B)
DEPTH_600 = RGBColor(0x1D, 0x64, 0x73)
DEPTH_100 = RGBColor(0xD3, 0xE7, 0xEA)
SAND = RGBColor(0xC9, 0xB7, 0x9A)
EARTH = RGBColor(0xA0, 0x8A, 0x6D)
CRITICAL = RGBColor(0xB9, 0x1C, 0x1C)
SAFE = RGBColor(0x15, 0x80, 0x3D)
INK = RGBColor(0x1C, 0x1A, 0x17)
INK_SOFT = RGBColor(0x5C, 0x57, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xF4, 0xF2, 0xEE)

BODY = "Calibri"
HEAD = "Cambria"

# ---------------------------------------------------------------- helpers ---


def kill(shape) -> None:
    shape._element.getparent().remove(shape._element)


def textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def line(tf, text, *, size=10.5, bold=False, color=INK, font=BODY,
         space_before=0, space_after=2, first=False, indent=0, align=None):
    """Append a paragraph. `first` reuses the empty paragraph a text frame starts with."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if indent:
        p.level = indent
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return p


def pointer(tf, text, *, first=False):
    """A template pointer, reused verbatim as a sub-heading."""
    return line(tf, text, size=12, bold=True, color=DEPTH_700, font=HEAD,
                space_before=0 if first else 9, space_after=4, first=first)


def bullet(tf, text, *, size=10.5, color=INK):
    return line(tf, "•  " + text, size=size, color=color, space_after=3)


def box(slide, shape, x, y, w, h, fill, *, outline=None, text=None,
        size=10, bold=True, color=WHITE, font=BODY, radius=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if outline is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = outline
        s.line.width = Pt(1)
    s.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = radius
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if text is not None:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for i, part in enumerate(text.split("\n")):
            if i:
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = part
            r.font.size = Pt(size if i == 0 else size - 1.5)
            r.font.bold = bold if i == 0 else False
            r.font.color.rgb = color
            r.font.name = font
    return s


def caption(slide, x, y, w, text, *, size=8.5, color=INK_SOFT, align=PP_ALIGN.LEFT,
            bold=False, h=None):
    # Height defaults to roughly the wrapped line count so QA can measure it;
    # PowerPoint grows a text box anyway, but a truthful box makes the check
    # meaningful.
    if h is None:
        # Calibri averages ~0.48 em per character, so w inches holds about
        # w * 150 / size characters. Line box is 1.25 x the point size.
        est_chars = max(int(w * 150 / size), 1)
        lines = max(1, -(-len(text) // est_chars))
        h = (size * 1.25 * lines) / 72 + 0.04
    tf = textbox(slide, x, y, w, h)
    line(tf, text, size=size, color=color, first=True, align=align, bold=bold)


def find(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle.lower() in sh.text_frame.text.lower():
            return sh
    return None


def pointer_lines(shape):
    """The template's own pointer strings, in order, blanks removed."""
    return [ln.strip() for ln in shape.text_frame.text.split("\n") if ln.strip()]


# ------------------------------------------------------------------ build ---

prs = Presentation(str(TEMPLATE))
S = prs.slides

IDEA_TITLE = "INGRES: groundwater answers you can check"
TEAM_NAME = "<< Team Name >>"


def set_title(slide, text):
    """Retitle a placeholder without losing its look.

    `text_frame.clear()` throws away the run properties the layout supplies, so
    a new run comes out as unstyled body text. Overwriting the existing run
    keeps the template's size, weight and colour.
    """
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == 0:
            for pi, p in enumerate(sh.text_frame.paragraphs):
                runs = list(p.runs)
                for ri, r in enumerate(runs):
                    r.text = text if (pi == 0 and ri == 0) else ""
            return


def set_team(slide, text):
    """Replace the template's team placeholder, keeping its formatting.

    Written by index, not identity: `p.runs` builds fresh proxy objects on
    every access, so `r is p.runs[0]` is never true and the original text
    survived untouched.
    """
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "Your Team Name":
            done = False
            for p in sh.text_frame.paragraphs:
                runs = list(p.runs)
                for i, r in enumerate(runs):
                    if i == 0 and not done:
                        r.text = text
                        done = True
                    else:
                        r.text = ""
            return

# ============================================================ 1 · title ======
s1 = S[0]
tb = find(s1, "Problem Statement ID")
tf = tb.text_frame
tf.clear()
rows = [
    ("Problem Statement ID –", "<< fill from portal >>"),
    ("Problem Statement Title –", "AI-driven assistant for groundwater data access"),
    ("Theme –", "<< fill from portal >>"),
    ("PS Category –", "Software"),
    ("Team ID –", "<< fill from portal >>"),
    ("Team Name –", "<< your registered team name >>"),
]
for i, (label, value) in enumerate(rows):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    r1 = p.add_run()
    r1.text = label + " "
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = DEPTH_900
    r1.font.name = HEAD
    r2 = p.add_run()
    r2.text = value
    r2.font.size = Pt(15)
    r2.font.color.rgb = INK
    r2.font.name = BODY

# ==================================================== 2 · proposed solution ==
s2 = S[1]
set_title(s2, IDEA_TITLE)
ptr = find(s2, "Proposed Solution")
pts = pointer_lines(ptr)          # section label + the three pointers
kill(ptr)

caption(s2, 0.5, 1.28, 8.0, pts[0], size=12.5, color=DEPTH_900, bold=True)

tf = textbox(s2, 0.5, 1.75, 6.35, 5.0)
pointer(tf, pts[1], first=True)
bullet(tf, "Ask in plain English; get a grounded answer with a chart, a map and a citation for every figure")
bullet(tf, "Figures come from SQL over 36,879 CGWB station readings — never from the language model")
bullet(tf, "Two sources by question type: the database for how much water, the CGWB report for quality and causes")

pointer(tf, pts[2])
bullet(tf, "The data exists but sits in PDFs and portals a non-specialist cannot query")
bullet(tf, "Removes the specialist barrier while every figure keeps its station name and date")
bullet(tf, "Refuses questions outside Punjab groundwater instead of guessing")

pointer(tf, pts[3])
bullet(tf, "Verification-first: eight deterministic checks can block an answer before it is shown")
bullet(tf, "Three of five pipeline stages use no language model at all")
bullet(tf, "Model-checking-model was tried and failed — it approved a Bathinda station as the source for a Ludhiana figure")
bullet(tf, "Runs fully offline on a free local model: no internet, no API key, no cost per query")

# --- diagram: why the obvious approach fails ---
DX, DW = 7.15, 5.65
box(s2, MSO_SHAPE.RECTANGLE, DX, 1.75, DW, 4.08, PAPER)
caption(s2, DX + 0.15, 1.9, DW - 0.3, "WHY THE OBVIOUS APPROACH FAILS",
        size=9, color=DEPTH_700, bold=True)

# path A — typical chatbot
caption(s2, DX + 0.15, 2.28, DW - 0.3, "A typical chatbot", size=9.5, color=INK, bold=True)
aw, ay, ah = 1.62, 2.62, 0.52
for i, (lbl, fill, col) in enumerate([
        ("Question", DEPTH_100, DEPTH_900),
        ("Language\nmodel", CRITICAL, WHITE),
        ("Answer", DEPTH_100, DEPTH_900)]):
    box(s2, MSO_SHAPE.CHEVRON, DX + 0.15 + i * (aw + 0.03), ay, aw, ah,
        fill, text=lbl, size=9, color=col)
caption(s2, DX + 0.15, ay + ah + 0.06, DW - 0.3,
        "The number comes OUT of the model. Nothing can check it.",
        size=8.5, color=CRITICAL)

# path B — INGRES
caption(s2, DX + 0.15, 3.62, DW - 0.3, "INGRES", size=9.5, color=INK, bold=True)
bw, by, bh = 1.22, 3.96, 0.52
for i, (lbl, fill, col) in enumerate([
        ("Question", DEPTH_100, DEPTH_900),
        ("SQL", DEPTH_700, WHITE),
        ("Model\nwrites prose", SAND, DEPTH_900),
        ("8 checks", SAFE, WHITE)]):
    box(s2, MSO_SHAPE.CHEVRON, DX + 0.15 + i * (bw + 0.02), by, bw, bh,
        fill, text=lbl, size=8.5, color=col)
caption(s2, DX + 0.15, by + bh + 0.06, DW - 0.3,
        "The number goes INTO the model, and is checked again on the way out.",
        size=8.5, color=SAFE)

caption(s2, DX + 0.15, 5.42, DW - 0.3,
        "An unverifiable answer is shown as raw data, never as confident prose.",
        size=9, color=DEPTH_900, bold=True)

# ==================================================== 3 · technical approach ==
s3 = S[2]
ptr = find(s3, "Technologies to be used")
pts = pointer_lines(ptr)
kill(ptr)

tf = textbox(s3, 0.5, 1.3, 12.35, 0.4)
pointer(tf, pts[0], first=True)
tf2 = textbox(s3, 0.5, 1.72, 6.1, 1.2)
bullet(tf2, "Backend: FastAPI · PostgreSQL · SQLAlchemy · NumPy", size=10)
bullet(tf2, "Frontend: React · Recharts · Leaflet", size=10)
tf3 = textbox(s3, 6.75, 1.72, 6.1, 1.2)
bullet(tf3, "AI: Ollama running qwen2.5:7b locally · nomic-embed-text", size=10)
bullet(tf3, "No vector database — 178 passages is a NumPy dot product", size=10)

tf4 = textbox(s3, 0.5, 2.98, 12.35, 0.4)
pointer(tf4, pts[1], first=True)

# --- diagram: the pipeline ---
py, ph, gap = 3.5, 0.95, 0.04
stages = [
    ("Query\nunderstanding", "LLM", SAND, DEPTH_900),
    ("Retrieval", "code", DEPTH_700, WHITE),
    ("Calculation", "code", DEPTH_700, WHITE),
    ("Response", "LLM", SAND, DEPTH_900),
    ("Verification", "code", DEPTH_700, WHITE),
]
pw = (12.35 - gap * (len(stages) - 1)) / len(stages)
for i, (lbl, tag, fill, col) in enumerate(stages):
    x = 0.5 + i * (pw + gap)
    box(s3, MSO_SHAPE.CHEVRON, x, py, pw, ph, fill, text=lbl, size=11, color=col)
    caption(s3, x, py + ph + 0.05, pw, tag.upper(), size=8.5,
            color=DEPTH_600 if tag == "code" else EARTH, align=PP_ALIGN.CENTER,
            bold=True)

caption(s3, 0.5, py + ph + 0.32, 12.35,
        "Three of the five stages involve no language model. Every number passes through code.",
        size=10, color=DEPTH_900, bold=True)

# --- data provenance strip ---
sy = 5.18
box(s3, MSO_SHAPE.RECTANGLE, 0.5, sy, 12.35, 1.32, PAPER)
caption(s3, 0.65, sy + 0.1, 6.0, "WHERE THE DATA COMES FROM", size=8.5,
        color=DEPTH_700, bold=True)
tfp = textbox(s3, 0.65, sy + 0.36, 5.9, 0.9)
line(tfp, "CGWB CSV exports  →  cleaning  →  PostgreSQL", size=10, bold=True,
     color=DEPTH_900, first=True)
line(tfp, "36,879 readings · 1,607 stations · 1996–2024", size=9.5, color=INK_SOFT)
tfq = textbox(s3, 6.85, sy + 0.36, 5.9, 0.9)
line(tfq, "CGWB report PDF  →  chunking  →  vectors", size=10, bold=True,
     color=DEPTH_900, first=True)
line(tfq, "178 passages · printed pages 1–78 · 153 assessment blocks",
     size=9.5, color=INK_SOFT)

# ================================================ 4 · feasibility & viability ==
s4 = S[3]
ptr = find(s4, "feasibility of the idea")
pts = pointer_lines(ptr)
kill(ptr)

tf = textbox(s4, 0.5, 1.3, 7.55, 5.3)
pointer(tf, pts[0], first=True)
bullet(tf, "Built and running end to end on real CGWB data — a working prototype, not a concept")
bullet(tf, "Runs on one machine with no runtime network dependency; zero recurring cost")
bullet(tf, "Source data is public and officially published — no access or licensing risk")
bullet(tf, "133 automated tests cover the deterministic core")

pointer(tf, pts[1])
bullet(tf, "A language model inventing a plausible figure — the central risk for any water tool")
bullet(tf, "Single-user: one local model, no request queue")
bullet(tf, "One report in the semantic corpus; database is CGWB's published set to 2024, not live")
bullet(tf, "Deployment blocked in development — the network blocked every database port")

pointer(tf, pts[2])
bullet(tf, "Eight blocking checks plus an advisory model reviewer; failure degrades to raw data")
bullet(tf, "Add a request queue and per-question timeout for concurrent users")
bullet(tf, "Corpus expansion is metadata-ready — each passage records its section, districts and scope")
bullet(tf, "Deploy inside the cloud network, where the port restriction does not apply")

# --- diagram: the verification gate ---
GX, GW = 8.35, 4.5
box(s4, MSO_SHAPE.RECTANGLE, GX, 1.3, GW, 5.15, PAPER)
caption(s4, GX + 0.15, 1.45, GW - 0.3, "THE VERIFICATION GATE", size=9,
        color=DEPTH_700, bold=True)

box(s4, MSO_SHAPE.ROUNDED_RECTANGLE, GX + 0.9, 1.85, 2.7, 0.5, DEPTH_100,
    text="Draft answer", size=10, color=DEPTH_900, radius=0.25)
box(s4, MSO_SHAPE.DOWN_ARROW, GX + 2.1, 2.4, 0.3, 0.3, DEPTH_600)
box(s4, MSO_SHAPE.ROUNDED_RECTANGLE, GX + 0.9, 2.75, 2.7, 0.62, DEPTH_700,
    text="Eight checks\ncitations · figures · districts · years · units",
    size=10, color=WHITE, radius=0.25)

box(s4, MSO_SHAPE.DOWN_ARROW, GX + 1.05, 3.43, 0.28, 0.3, SAFE)
box(s4, MSO_SHAPE.DOWN_ARROW, GX + 3.15, 3.43, 0.28, 0.3, CRITICAL)
caption(s4, GX + 0.12, 3.44, 0.88, "pass", size=8.5, color=SAFE, bold=True,
        align=PP_ALIGN.RIGHT)
caption(s4, GX + 3.5, 3.44, 0.9, "fail", size=8.5, color=CRITICAL, bold=True)

box(s4, MSO_SHAPE.ROUNDED_RECTANGLE, GX + 0.2, 3.8, 2.0, 0.75, SAFE,
    text="Answer shown\nwith citations", size=9.5, color=WHITE, radius=0.2)
box(s4, MSO_SHAPE.ROUNDED_RECTANGLE, GX + 2.3, 3.8, 2.0, 0.75, CRITICAL,
    text="Raw data shown,\nmarked unverified", size=9.5, color=WHITE, radius=0.2)

caption(s4, GX + 0.15, 4.72, GW - 0.3,
        "An unverifiable answer is never shown as if it were verified. The system "
        "fails to data, not to silence and not to a confident guess.",
        size=9, color=DEPTH_900)

caption(s4, GX + 0.15, 5.45, GW - 0.3, "Six distinct failure paths are handled: "
        "out of scope · no quality data · district with no readings · checks "
        "rejected · model unusable · model unreachable.", size=8.5, color=INK_SOFT)

# ================================================== 5 · impact and benefits ===
s5 = S[4]
ptr = find(s5, "target audience")
pts = pointer_lines(ptr)
kill(ptr)

tf = textbox(s5, 0.5, 1.3, 6.6, 5.3)
pointer(tf, pts[0], first=True)
bullet(tf, "Farmers and panchayats deciding where and how deep to drill")
bullet(tf, "District and block officials, and CGWB field staff")
bullet(tf, "Policy researchers, students and journalists")

pointer(tf, pts[1])
line(tf, "Social", size=10.5, bold=True, color=DEPTH_900, space_before=5, space_after=2)
bullet(tf, "Groundwater status becomes answerable by anyone who can type a question")
bullet(tf, "Every answer is citable, so decisions taken from it can be audited")
line(tf, "Economic", size=10.5, bold=True, color=DEPTH_900, space_before=5, space_after=2)
bullet(tf, "Better-informed borewell and cropping decisions; avoids drilling into a falling table")
bullet(tf, "No licensing or per-query cost — deployable without a budget line")
line(tf, "Environmental", size=10.5, bold=True, color=DEPTH_900, space_before=5, space_after=2)
bullet(tf, "Makes a slow crisis visible: 20 of 23 assessed districts are over-exploited")
bullet(tf, "Trend and projection views turn depletion into something a non-specialist can see")

# --- map: district markers over the monitoring network's own outline ---
data = json.loads((HERE / "districts.json").read_text())
proj = json.loads((HERE / "punjab.json").read_text())

MX, MY, MW, MH = 7.4, 1.4, 5.4, 5.05
box(s5, MSO_SHAPE.RECTANGLE, MX, MY, MW, MH, PAPER)
caption(s5, MX + 0.18, MY + 0.13, MW - 0.36,
        "PUNJAB, BY CGWB CATEGORY", size=9, color=DEPTH_700, bold=True)

# Sized by height, not width: cropped to the landmass the map is portrait,
# and fitting it to the panel's width ran it through the footer.
IMG_Y, IMG_H = 1.76, 3.62
IMG_W = IMG_H * proj["w"] / proj["h"]
IMG_X = MX + (MW - IMG_W) / 2
s5.shapes.add_picture(str(HERE / "punjab.png"), Inches(IMG_X), Inches(IMG_Y),
                      Inches(IMG_W), Inches(IMG_H))


def place(lat, lon):
    """District coordinate -> slide inches, using the map's own projection."""
    fx = (proj["ox"] + (lon - proj["lon0"]) * proj["kx"] * proj["scale"]) / proj["w"]
    fy = (proj["oy"] + (proj["lat1"] - lat) * proj["scale"]) / proj["h"]
    return IMG_X + fx * IMG_W, IMG_Y + fy * IMG_H


COLOR = {"over-exploited": CRITICAL, "safe": SAFE}
dot = 0.155
for p in data["points"]:
    cx, cy = place(p["lat"], p["lon"])
    box(s5, MSO_SHAPE.OVAL, cx - dot / 2, cy - dot / 2, dot, dot,
        COLOR.get(p["c"], EARTH), outline=WHITE)

# Name the three exceptions. Everything else is the same colour and the same
# story, and 23 labels would bury the map.
for p in data["points"]:
    if p["c"] != "safe":
        continue
    cx, cy = place(p["lat"], p["lon"])
    caption(s5, cx - 0.72, cy + 0.07, 1.5, p["d"], size=7, color=DEPTH_900,
            bold=True, align=PP_ALIGN.CENTER)

ly = MY + MH - 0.92
for i, (lbl, col, n) in enumerate([("Over-exploited — 20", CRITICAL, 20),
                                   ("Safe — 3", SAFE, 3)]):
    box(s5, MSO_SHAPE.OVAL, MX + 0.3 + i * 2.5, ly, 0.13, 0.13, col)
    caption(s5, MX + 0.5 + i * 2.5, ly - 0.035, 2.2, lbl, size=9, color=INK)
box(s5, MSO_SHAPE.OVAL, MX + 0.3, ly + 0.26, 0.09, 0.09,
    RGBColor(0x60, 0x7A, 0x80))
caption(s5, MX + 0.5, ly + 0.21, 4.6,
        "1,606 monitoring stations — the outline is their own coverage",
        size=8.5, color=INK_SOFT)

caption(s5, MX + 0.3, ly + 0.53, MW - 0.6,
        "Malerkotla is categorised by CGWB but has no monitoring stations, so it "
        "is named here rather than plotted.", size=7.5, color=INK_SOFT)

# ================================================ 6 · research & references ===
s6 = S[5]
ptr = find(s6, "reference and research")
pts = pointer_lines(ptr)
kill(ptr)

tf = textbox(s6, 0.5, 1.35, 12.35, 5.2)
pointer(tf, pts[0], first=True)

refs = [
    ("CGWB — Ground Water Resources of Punjab (as on 31 March 2024)",
     "Block-wise resources and the water-quality chapters. Source of all 153 block "
     "assessments and the 178 indexed passages."),
    ("National Water Data Portal (NWIC) — Ground Water Level, Manual Quarterly, CGWB",
     "36,879 readings across 1,607 stations, 1996–2024.  nwdp.nwic.gov.in"),
    ("GEC-2015 — Ground Water Estimation Committee methodology",
     "The assessment methodology CGWB's resource estimation follows."),
    ("BIS IS 10500 — Drinking Water Specification",
     "Limits used in the report's quality chapter, e.g. fluoride 1.5 mg/L, uranium 0.03 mg/L."),
    ("CGWB INGRES — India Ground Water Resource Estimation System",
     "The national assessment system this assistant is designed to support.  ingres.iith.ac.in"),
]
for title, detail in refs:
    line(tf, "•  " + title, size=11, bold=True, color=DEPTH_900, space_before=7,
         space_after=1)
    line(tf, "    " + detail, size=9.5, color=INK_SOFT, space_after=0)

line(tf, "Methodology note", size=10.5, bold=True, color=DEPTH_700, font=HEAD,
     space_before=12, space_after=2)
line(tf, "CGWB classifies assessment blocks, never districts. District categories "
         "shown by the system are derived from block counts, which are stored "
         "alongside as the supporting evidence. Cross-checking the extraction "
         "against a second table in the same report corrected four categories.",
     size=9.5, color=INK)

# ================================================ drop the instructions page ==
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
rId = ids[6].get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
prs.part.drop_rel(rId)
sldIdLst.remove(ids[6])

for _s in list(S)[1:]:
    set_team(_s, TEAM_NAME)
    # The template sizes this box for a short name; widen it so a real team
    # name sits on one line instead of wrapping under itself.
    for _sh in _s.shapes:
        if _sh.has_text_frame and _sh.text_frame.text.strip() == TEAM_NAME:
            _sh.width = Inches(3.0)

prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
